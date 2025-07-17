# app.py

import streamlit as st
import pandas as pd
import uuid
import json
from fpdf import FPDF
from datetime import datetime
import re
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import random

# This block runs first to handle setting the newly created session ID
# before the selectbox widget is created, preventing the "cannot be modified" error.
if 'just_created_session_id' in st.session_state:
    st.session_state.selected_session_id = st.session_state.just_created_session_id
    del st.session_state.just_created_session_id


# ======================================================================================
# Session State Initialization
# ======================================================================================

if 'questions' not in st.session_state: st.session_state.questions = {}
if 'company_info' not in st.session_state: st.session_state.company_info = {}
if 'company_summary' not in st.session_state: st.session_state.company_summary = ""
if 'roadmap_df' not in st.session_state: st.session_state.roadmap_df = pd.DataFrame()
if 'messages' not in st.session_state: st.session_state.messages = []
if 'notes_content' not in st.session_state: st.session_state.notes_content = ""
if 'value_strategy_content' not in st.session_state: st.session_state.value_strategy_content = ""
if 'competitive_analysis_content' not in st.session_state: st.session_state.competitive_analysis_content = ""
if 'session_loaded' not in st.session_state: st.session_state.session_loaded = False
if 'selected_session_id' not in st.session_state: st.session_state.selected_session_id = "new"
if 'selected_model' not in st.session_state: st.session_state.selected_model = 'claude-3-5-sonnet'
if 'outreach_emails' not in st.session_state: st.session_state.outreach_emails = {}
if 'initial_value_hypothesis' not in st.session_state: st.session_state.initial_value_hypothesis = ""
if 'research_stage' not in st.session_state: st.session_state.research_stage = 0
if 'recommended_initiatives' not in st.session_state: st.session_state.recommended_initiatives = []
if 'token_error_flag' not in st.session_state: st.session_state.token_error_flag = False
if 'people_research' not in st.session_state: st.session_state.people_research = []
if 'show_briefing' not in st.session_state: st.session_state.show_briefing = False

# ======================================================================================
# App Configuration
# ======================================================================================

st.set_page_config(
    page_title="Snowflake Sales Discovery Assistant",
    page_icon="❄️",
    layout="wide",
)

st.title("❄️ Snowflake Sales Discovery Assistant")
st.caption("An internal app for account research and discovery powered by Streamlit in Snowflake.")


# ======================================================================================
# Snowflake Connection and Data Functions
# ======================================================================================

conn = st.connection("snowflake")

def get_current_user():
    if 'current_user' not in st.session_state:
        try:
            user_df = conn.query("SELECT CURRENT_USER() as user;")
            st.session_state.current_user = user_df.iloc[0]['USER']
        except Exception as e:
            st.error(f"Could not retrieve current user: {e}")
            st.session_state.current_user = "UNKNOWN_USER"
    return st.session_state.current_user


def save_answers_to_snowflake(df):
    if df.empty:
        st.warning("No answered questions to save.")
        return False
        
    try:
        df_to_save = df.copy()
        df_to_save.columns = [col.upper() for col in df_to_save.columns]
        conn.write_pandas(df_to_save, "SALES_DISCOVERY_ANSWERS")
        return True
    except Exception as e:
        st.error(f"Error saving to Snowflake: {e}")
        return False

def get_next_session_version(session_name_prefix, user):
    sql = f"SELECT MAX(REGEXP_SUBSTR(session_name, 'v(\\\\d+)$')) FROM SALES_DISCOVERY_SESSIONS WHERE session_name LIKE '{session_name_prefix}%' AND saved_by_user = '{user}'"
    try:
        cursor = conn.cursor()
        cursor.execute(sql)
        result = cursor.fetchone()
        if result and result[0]:
            version_str = result[0]
            numeric_part = re.sub(r'\D', '', version_str)
            return int(numeric_part) + 1
        return 1
    except Exception as e:
        st.error(f"Error getting next session version: {e}")
        return 1

def save_session_to_snowflake(show_message=True):
    if not st.session_state.company_info.get('website'):
        if show_message:
            st.warning("Cannot save an empty session. Please generate questions for a company first.")
        return

    current_user = get_current_user()
    
    roadmap_json = st.session_state.roadmap_df.to_json(orient='split') if not st.session_state.roadmap_df.empty else None

    session_state_data = {
        "company_info": st.session_state.company_info,
        "questions": st.session_state.questions,
        "company_summary": st.session_state.company_summary,
        "notes_content": st.session_state.notes_content,
        "saved_by_user": current_user,
        "value_strategy_content": st.session_state.get('value_strategy_content', ''),
        "competitive_analysis_content": st.session_state.get('competitive_analysis_content', ''),
        "roadmap_json": roadmap_json,
        "selected_model": st.session_state.get('selected_model', 'claude-3-5-sonnet'),
        "outreach_emails": st.session_state.get('outreach_emails', {}),
        "initial_value_hypothesis": st.session_state.get('initial_value_hypothesis', ''),
        "research_stage": st.session_state.get('research_stage', 0),
        "recommended_initiatives": st.session_state.get('recommended_initiatives', []),
        "people_research": st.session_state.get('people_research', [])
    }
    session_state_json = json.dumps(session_state_data)

    if st.session_state.selected_session_id != "new":
        try:
            session_id = st.session_state.selected_session_id
            sql = "UPDATE SALES_DISCOVERY_SESSIONS SET session_state = PARSE_JSON(?) WHERE session_id = ? AND saved_by_user = ?"
            with conn.cursor() as cursor:
                cursor.execute(sql, (session_state_json, session_id, current_user))
                if cursor.rowcount > 0:
                    if show_message:
                        st.success("Session updated successfully!")
                else:
                    st.warning("Could not update. Session not found or you don't have permission.")
        except Exception as e:
            st.error(f"Failed to update session: {e}")
    
    else:
        company_name = st.session_state.company_info['website'].replace('https://','').replace('www.','').split('.')[0]
        date_str = datetime.now().strftime('%Y-%m-%d')
        session_name_prefix = f"{company_name}_{date_str}"
        version = get_next_session_version(session_name_prefix, current_user)
        
        random_suffix = random.randint(100, 999)
        session_name = f"{session_name_prefix}_v{version}_{random_suffix}"
        session_id = str(uuid.uuid4())
        
        sql = "INSERT INTO SALES_DISCOVERY_SESSIONS (session_id, session_name, session_state, saved_by_user) SELECT ?, ?, PARSE_JSON(?), ?"
        try:
            with conn.cursor() as cursor:
                cursor.execute(sql, (session_id, session_name, session_state_json, current_user))
            
            if show_message:
                st.success(f"Session '{session_name}' saved successfully!")
            
            st.session_state.just_created_session_id = session_id
            st.session_state.session_loaded = True
        except Exception as e:
            st.error(f"Failed to save new session: {e}")


def load_sessions_from_snowflake():
    current_user = get_current_user()
    try:
        query = "SELECT session_id, session_name FROM SALES_DISCOVERY_SESSIONS WHERE saved_by_user = ? ORDER BY created_at DESC;"
        df = conn.query(query, params=(current_user,), ttl=0)
        return df
    except Exception as e:
        if "does not exist" in str(e):
            st.sidebar.warning("SALES_DISCOVERY_SESSIONS table not found. Creating it now. Please refresh.")
            create_sessions_table_sql = """
            CREATE OR REPLACE TABLE SALES_DISCOVERY_SESSIONS (
                session_id VARCHAR PRIMARY KEY, session_name VARCHAR,
                created_at TIMESTAMP_NTZ DEFAULT CURRENT_TIMESTAMP(),
                session_state VARIANT, saved_by_user VARCHAR
            );
            """
            with conn.cursor() as cursor:
                cursor.execute(create_sessions_table_sql)
            return pd.DataFrame(columns=['SESSION_ID', 'SESSION_NAME'])
        else:
            st.error(f"Error loading sessions: {e}")
            return pd.DataFrame(columns=['SESSION_ID', 'SESSION_NAME'])


def load_selected_session(session_id):
    current_user = get_current_user()
    try:
        query = "SELECT session_state FROM SALES_DISCOVERY_SESSIONS WHERE session_id = ? AND saved_by_user = ?;"
        df = conn.query(query, params=(session_id, current_user), ttl=0)
        if not df.empty:
            session_data_str = df.iloc[0]['SESSION_STATE']
            session_data = json.loads(session_data_str)
            
            clear_session(preserve_loaded_session_id=True)
            
            st.session_state.update(session_data)
            st.session_state.roadmap_df = pd.read_json(session_data.get('roadmap_json', '{}'), orient='split') if session_data.get('roadmap_json') else pd.DataFrame()
            st.session_state.research_stage = session_data.get('research_stage', 2)
            st.session_state.session_loaded = True
        else:
            st.warning("Could not find the selected session or you do not have permission to access it.")
    except Exception as e:
        st.error(f"Error loading session state: {e}")

def delete_session_from_snowflake(session_id):
    current_user = get_current_user()
    sql = "DELETE FROM SALES_DISCOVERY_SESSIONS WHERE session_id = ? AND saved_by_user = ?"
    try:
        with conn.cursor() as cursor:
            cursor.execute(sql, (session_id, current_user))
            if cursor.rowcount > 0:
                st.success("Session deleted.")
            else:
                st.warning("Session not found or you don't have permission to delete it.")
        st.session_state.selected_session_id = "new" 
    except Exception as e:
        st.error(f"Failed to delete session: {e}")

# ======================================================================================
# LLM Functions
# ======================================================================================

def cortex_request(prompt, json_output=True):
    selected_model = st.session_state.get('selected_model', 'claude-3-5-sonnet')
    
    llm_response_str = ""
    try:
        safe_prompt = prompt.replace("'", "''")
        sql_query = f"SELECT SNOWFLAKE.CORTEX.COMPLETE('{selected_model}', '{safe_prompt}') as response;"
        cursor = conn.cursor()
        cursor.execute(sql_query)
        response_row = cursor.fetchone()
        if not (response_row and response_row[0]):
            st.error("The AI model returned an empty response.")
            return None
        llm_response_str = response_row[0]
    except Exception as e:
        if "max tokens" in str(e) and "exceeded" in str(e):
            st.session_state.token_error_flag = True
        else:
            st.error(f"An error occurred while calling the LLM: {e}")
        return None

    if not json_output:
        return llm_response_str

    try:
        json_start = llm_response_str.find('{')
        json_end = llm_response_str.rfind('}') + 1
        if json_start != -1 and json_end != 0:
            json_str = llm_response_str[json_start:json_end]
            return json.loads(json_str)
        else:
            raise json.JSONDecodeError("No JSON object found in response", llm_response_str, 0)
    except json.JSONDecodeError:
        st.warning("Initial JSON parsing failed. Attempting to repair with another LLM call...")
        
        repair_prompt_text = f"The following text is not valid JSON. Please correct the syntax and return only the corrected, valid JSON object. Do not add any commentary or explanation. Here is the invalid text:\n\n{llm_response_str}"
        safe_repair_prompt = repair_prompt_text.replace("'", "''")
        
        try:
            repair_sql_query = f"SELECT SNOWFLAKE.CORTEX.COMPLETE('{selected_model}', '{safe_repair_prompt}') as response;"
            cursor = conn.cursor()
            cursor.execute(repair_sql_query)
            repair_response_row = cursor.fetchone()
            if not (repair_response_row and repair_response_row[0]):
                st.error("LLM repair attempt returned an empty response.")
                return None
            
            repaired_llm_str = repair_response_row[0]
            json_start = repaired_llm_str.find('{')
            json_end = repaired_llm_str.rfind('}') + 1

            if json_start != -1 and json_end != 0:
                repaired_json_str = repaired_llm_str[json_start:json_end]
                return json.loads(repaired_json_str)
            else:
                st.error("LLM JSON repair failed. Could not parse the response.")
                st.text_area("Original LLM response that failed parsing", llm_response_str, height=200)
                return None
        except Exception as e:
            st.error(f"An error occurred during the LLM repair attempt: {e}")
            st.text_area("Original LLM response that failed parsing", llm_response_str, height=200)
            return None

def research_person(name, title, company_info):
    """Generates conversation starters for a specific person."""
    company_name = company_info.get('website', 'their company')
    prompt = f"""
    You are a professional relationship-building assistant. Your goal is to find insightful and relevant conversation starters.
    Based on your public knowledge about a person named **{name}**, who holds the title of **{title}** at **{company_name}**, generate potential topics for discussion.

    Return your response as a single, valid JSON object with the following keys:
    - "summary": A concise, one-paragraph summary of this person's likely professional focus and expertise, based on their title and company.
    - "insights": A list of 3-5 bullet points identifying key technical skills (e.g., Python, SQL, AWS, AI/ML) and business expertise (e.g., e-commerce, risk management, supply chain) this person might have.
    - "topics": A JSON object with three keys: "business", "technical", and "personal". Each key should contain a list of 2-3 suggested questions or topics you could discuss with them to build rapport. For "personal", suggest topics based on common interests for someone in their role and industry (e.g., recent industry conferences, popular tech books, local events).
    """
    return cortex_request(prompt)

def generate_discovery_questions(website, industry, competitor, persona):
    prompt = f"""
    You are an expert Snowflake sales engineer. Generate discovery questions for a potential customer with the title of **{persona}**.
    Company Information: Website: {website}, Industry: {industry}, Assumed Primary Competitor: {competitor}.
    Generate three categories: **Technical Discovery**, **Business Discovery**, and exactly 10 questions for **Competitive Positioning vs. {competitor}**.
    The total number of questions must not exceed 30. Return as a single, valid JSON object.
    """
    questions_dict = cortex_request(prompt)
    if questions_dict:
        temp_processed = {}
        competitive_key_name = f"Competitive Positioning vs. {competitor}"
        for category, qs in questions_dict.items():
            if "Competitive Positioning" in category or competitor in category:
                temp_processed[competitive_key_name] = [{"id": str(uuid.uuid4()), "text": q, "answer": "", "favorite": False} for q in qs]
            else:
                temp_processed[category.title()] = [{"id": str(uuid.uuid4()), "text": q, "answer": "", "favorite": False} for q in qs]
        
        ordered_questions = {}
        category_order = ["Technical Discovery", "Business Discovery", competitive_key_name]

        for category in category_order:
            if category.title() in temp_processed:
                ordered_questions[category.title()] = temp_processed[category.title()]
            elif category in temp_processed:
                ordered_questions[category] = temp_processed[category]
        
        for category, questions in temp_processed.items():
            if category not in ordered_questions:
                ordered_questions[category] = questions
                
        return ordered_questions
    return {}

def generate_company_summary(website, industry, persona):
    prompt = f"""
    You are a skilled business analyst. Based on the company website '{website}' and its industry '{industry}', provide a concise, one-paragraph summary relevant to a person with the title of **{persona}**.
    Touch upon the company's likely business model, its target customers, and potential data-related opportunities or challenges.
    """
    return cortex_request(prompt, json_output=False)

def recommend_key_initiatives(company_info):
    initiatives_list = ["Grow Revenue", "Reduce Risk", "Manage Costs", "Innovate"]
    prompt = f"""
    You are a business strategy analyst. Based on the company with website '{company_info.get('website')}' in the '{company_info.get('industry')}' industry, what are their three most likely key business initiatives?
    Choose from the following standard list: {initiatives_list}.
    If you believe a more specific initiative is highly likely, you may suggest one custom initiative, but it must be concise (3-5 words).
    Your response MUST be a single, valid JSON object with a key "recommended_initiatives" which contains a list of exactly 3 English strings.
    """
    response = cortex_request(prompt)
    if response and "recommended_initiatives" in response and isinstance(response["recommended_initiatives"], list):
        return response["recommended_initiatives"]
    return []

def generate_more_questions_for_category(website, industry, competitor, persona, category, existing_questions):
    existing_questions_str = "\n".join([f"- {q}" for q in existing_questions])
    prompt = f"""
    You are an expert Snowflake sales engineer. For a **{persona}** at a company with website {website}, generate exactly 5 new and distinct discovery questions for the category '{category}'.
    Do not repeat these existing questions: {existing_questions_str}.
    Return your response as a single, valid JSON object with a key "new_questions" containing a list of strings.
    """
    response_dict = cortex_request(prompt)
    if response_dict and "new_questions" in response_dict:
        return [{"id": str(uuid.uuid4()), "text": q, "answer": "", "favorite": False} for q in response_dict["new_questions"]]
    return []

def generate_initiative_questions(website, industry, persona, initiative):
    prompt = f"""
    You are an expert Snowflake sales engineer. For a **{persona}** at a company with website {website}, generate 5-7 insightful discovery questions that directly relate to the business initiative: **{initiative}**.
    Return your response as a single, valid JSON object with a key "questions" containing a list of strings.
    """
    response_dict = cortex_request(prompt)
    if response_dict and "questions" in response_dict:
        return [{"id": str(uuid.uuid4()), "text": q, "answer": "", "favorite": False} for q in response_dict["questions"]]
    return []

def generate_roadmap(company_info, discovery_notes_str, priority):
    prompt = f"""
    You are a world-class Snowflake Solution Architect creating a strategic roadmap for a potential customer.
    **Customer Context:** Website: {company_info.get('website', 'N/A')}, Industry: {company_info.get('industry', 'N/A')}, Persona: {company_info.get('persona', 'N/A')}.
    **Discovery Notes:**
    {discovery_notes_str}
    **Roadmap Requirements:**
    1. Analyze the notes to identify key pain points and business objectives.
    2. Propose a logical sequence of 2-4 implementation projects using Snowflake.
    3. For each project, provide a `project_name`, `description`, `business_value` ('Low', 'Medium', 'High', 'Very High'), and `level_of_effort` ('Low', 'Medium', 'High').
    4. Order the roadmap based on this priority: **{priority}**.
    5. You MUST return the output as a single, valid JSON object with a key "roadmap" which is a list of project objects.
    """
    response_dict = cortex_request(prompt)
    if response_dict and "roadmap" in response_dict:
        return pd.DataFrame(response_dict["roadmap"])
    return pd.DataFrame()

def get_chatbot_response(company_info, chat_history):
    history_str = ""
    for message in chat_history:
        role = "User" if message['role'] == 'user' else "Assistant"
        history_str += f"{role}: {message['content']}\n"
    prompt = f"""
    You are a helpful Snowflake sales assistant chatbot. Your goal is to answer user questions about how Snowflake's technology can help a specific company.
    You have the following context:
    - Company Website: {company_info.get('website', 'N/A')}
    - Industry: {company_info.get('industry', 'N/A')}
    - Assumed Primary Competitor: {company_info.get('competitor', 'N/A')}
    - Contact's Title: {company_info.get('persona', 'N/A')}
    You also have the conversation history:
    {history_str}
    Based on all of this context, provide a concise and helpful answer to the last user question.
    """
    return cortex_request(prompt, json_output=False)

def autofill_answers_from_notes(notes_text, questions_dict):
    populated_questions = {}
    for category, questions in questions_dict.items():
        if not questions:
            continue
        st.write(f"Analyzing notes for category: **{category}**...")
        answered_category = autofill_category_from_notes(notes_text, category, questions)
        if answered_category:
            populated_questions[category] = answered_category
        else:
            populated_questions[category] = questions
            st.warning(f"Could not auto-fill answers for category: {category}. Please review manually.")
    return populated_questions

def autofill_category_from_notes(notes_text, category, questions_in_category):
    max_chars = 12000 
    if len(notes_text) > max_chars:
        notes_text = notes_text[:max_chars]
    
    prompt = f"""
    You are an intelligent assistant. Your task is to populate answers for a single category of discovery questions based on a provided block of freeform notes.
    **Source Notes:**
    ---
    {notes_text}
    ---
    **Questions for the '{category}' category to Answer (JSON format):**
    ---
    {json.dumps(questions_in_category, indent=2)}
    ---
    **Instructions:**
    1. Read through the **Source Notes** carefully.
    2. For each question in the provided JSON object, find the relevant information within the notes.
    3. Formulate a concise answer for each question based *only* on the information present in the notes.
    4. If no relevant information for a question is found, the value for the "answer" key for that question MUST be an empty string ("").
    5. Your final output MUST be a single, valid JSON object with a key named "answered_questions" that contains a list of the questions with their "answer" fields populated. The structure of each question object must be identical to the input.
    6. CRITICAL: The returned JSON must be perfectly formatted. Pay very close attention to adding a comma `,` between each object in the list. Do not include a trailing comma after the last object in the list.
    """
    response = cortex_request(prompt)
    if response and "answered_questions" in response:
        return response["answered_questions"]
    return None

def generate_initial_value_hypothesis(company_info):
    prompt = f"""
    You are a senior business value consultant for Snowflake. Your task is to create an initial value hypothesis for a potential customer *before* a formal discovery call.
    **Customer Context:**
    - Company Website: {company_info.get('website', 'N/A')}
    - Industry: {company_info.get('industry', 'N/A')}
    **Instructions:**
    Based *only* on the publicly available information implied by the website and industry, generate a compelling value hypothesis. Structure your response in Markdown with the following sections:
    1.  **Executive Summary:** A brief overview of the likely value Snowflake can provide.
    2.  **Hypothesized Value Drivers:**
        * **Making Money:** How can this company likely leverage data with Snowflake to increase revenue or create new revenue streams?
        * **Reducing Costs:** What are the most probable areas where Snowflake could help reduce operational or infrastructure costs?
        * **Driving Innovation:** How might Snowflake enable them to innovate faster, for example, with AI/ML or new data applications?
    3.  **Potential Snowflake Use Cases & Suspected Value:**
        * Provide a list of 3-5 specific, high-impact use cases for Snowflake tailored to this company.
        * For each use case, provide a name, a brief explanation, and assign a 'Suspected Business Value' (e.g., Very High, High, Medium, Low). Format this as a list.
    This hypothesis is a starting point for a sales conversation. The entire response must be in well-formatted Markdown.
    """
    return cortex_request(prompt, json_output=False)

def generate_business_case(company_info, discovery_notes_str):
    prompt = f"""
    You are a senior business value consultant. Your task is to analyze the following discovery notes and build a compelling business case for adopting Snowflake.
    **Customer Context:**
    - Company: {company_info.get('website', 'N/A')}
    - Industry: {company_info.get('industry', 'N/A')}
    - Persona of Contact: {company_info.get('persona', 'N/A')}
    - Main Competitor: {company_info.get('competitor', 'N/A')}
    **Discovery Notes (Questions & Answers):**
    ---
    {discovery_notes_str}
    ---
    **Instructions:**
    Based on the notes, generate a response with three sections:
    1.  **Business Case:** Synthesize the notes into a business case. Use metrics and quantitative data from the answers where possible. Frame the value proposition around the key business initiatives discussed.
    2.  **Key Questions for Value Metrics:** Identify and list 3-5 critical follow-up questions that need to be asked to uncover stronger, quantifiable business value metrics.
    3.  **Recommended Strategy:** Suggest a high-level strategy to strengthen the business case and align with the customer's goals.
    Format your response in Markdown.
    """
    return cortex_request(prompt, json_output=False)

def generate_competitive_argument(company_info, discovery_notes_str):
    prompt = f"""
    You are a highly skilled, aggressive, and effective salesperson for **{company_info.get('competitor', 'the competitor')}**. 
    Your goal is to build the strongest possible "steel man" argument to convince a customer to choose your solution over Snowflake.
    **Customer Context:**
    - Company: {company_info.get('website', 'N/A')}
    - Industry: {company_info.get('industry', 'N/A')}
    - Persona of Contact: {company_info.get('persona', 'N/A')}
    **Your Competitor's Discovery Notes (Snowflake's perspective):**
    ---
    {discovery_notes_str}
    ---
    **Instructions:**
    Based on the notes, create the best possible competitive strategy and argument for **{company_info.get('competitor', 'your company')}**.
    1.  **Competitive Strategy:** What is your overall strategy to win this deal?
    2.  **Key Talking Points:** What are the 3-5 most powerful talking points you will use?
    3.  **How to Counter Snowflake:** How would you proactively counter Snowflake's main value propositions?
    Give the best, most compelling argument possible. Be specific and tactical. Format the response in Markdown.
    """
    return cortex_request(prompt, json_output=False)

def generate_outreach_emails(company_info, discovery_notes_str, roadmap_df):
    roadmap_str = roadmap_df.to_string() if not roadmap_df.empty else "No roadmap generated yet."
    prompt = f"""
    You are an expert enterprise software salesperson. Your goal is to draft 3 distinct, succinct, and captivating outreach emails to a person with the title of **{company_info.get('persona')}** at the company visited at **{company_info.get('website')}**.
    Use the following context to make the emails relevant and impactful:
    **Customer's Key Initiatives & Pains (from discovery notes):**
    ---
    {discovery_notes_str}
    ---
    **Potential Snowflake Projects We've Identified:**
    ---
    {roadmap_str}
    ---
    **Instructions:**
    1.  Draft three different versions of a follow-up email.
    2.  Each email should be short (less than 150 words).
    3.  Each email should reference a specific pain point from the discovery notes and tie it to a high-value or "quick win" (low-effort) project from the roadmap.
    4.  The tone should be professional, helpful, and focused on generating the next meeting.
    5.  Return the three email versions as a single, valid JSON object with keys "email_1", "email_2", and "email_3". Each key's value should be another JSON object containing "subject" and "body".
    6.  IMPORTANT: The final output must be only the JSON object. All strings within the JSON, especially the 'body', must be properly formatted with special characters like newlines escaped as \\n.
    """
    return cortex_request(prompt)

def regenerate_single_email(company_info, discovery_notes_str, roadmap_df, existing_emails, email_to_replace_key):
    other_emails_str = ""
    for key, data in existing_emails.items():
        if key != email_to_replace_key:
            other_emails_str += f"Existing Draft ({key}):\nSubject: {data.get('subject')}\nBody: {data.get('body')}\n\n"
    roadmap_str = roadmap_df.to_string() if not roadmap_df.empty else "No roadmap generated yet."
    prompt = f"""
    You are an expert salesperson. Redraft a single, captivating outreach email for a {company_info.get('persona')} at {company_info.get('website')}.
    Here is the original context:
    - Customer Pains: {discovery_notes_str}
    - Proposed Projects: {roadmap_str}
    Here are the other email drafts you have already written. DO NOT REPEAT OR REUSE the style or content from these existing drafts:
    ---
    {other_emails_str}
    ---
    INSTRUCTION:
    Generate one new, completely distinct email version that is different from the existing drafts. It must be short and professional.
    Return it as a single, valid JSON object with "subject" and "body" keys.
    IMPORTANT: The final output must be only the JSON object. All strings within the JSON must be properly formatted with special characters like newlines escaped as \\n.
    """
    return cortex_request(prompt)

# ======================================================================================
# Helper Functions for UI, State, and Exporting
# ======================================================================================

class PDF(FPDF):
    def __init__(self, orientation='P', unit='mm', format='A4', title=''):
        super().__init__(orientation, unit, format)
        self.title = title
    def header(self):
        self.set_font("Helvetica", "B", 12)
        self.cell(0, 10, self.title, 0, 0, 'C')
        self.ln(10)
    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

def create_notes_pdf_bytes(export_df, company_info):
    pdf = PDF(title='Snowflake Discovery Notes')
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    clean_website = company_info.get('website', 'N/A').encode('latin-1', 'replace').decode('latin-1')
    pdf.cell(0, 10, f"Discovery Notes: {clean_website}", 0, 1, 'C')
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 5, f"Industry: {company_info.get('industry', 'N/A')}", 0, 1)
    pdf.cell(0, 5, f"Persona: {company_info.get('persona', 'N/A')}", 0, 1)
    pdf.cell(0, 5, f"Assumed Competitor: {company_info.get('competitor', 'N/A')}", 0, 1)
    pdf.ln(10)
    for category, group in export_df.groupby('question_category'):
        pdf.set_font("Helvetica", "B", 12)
        clean_category = category.encode('latin-1', 'replace').decode('latin-1')
        pdf.cell(0, 10, clean_category, 0, 1, 'L', True)
        pdf.ln(2)
        for _, row in group.iterrows():
            fav_char = " (Favorite)" if row['favorite'] else ""
            question_text = f"Q: {row['question']}{fav_char}"
            answer_text = f"A: {row['answer'] if row['answer'] else 'No answer provided.'}"
            pdf.set_font("Helvetica", "B", 10)
            pdf.multi_cell(0, 5, question_text.encode('latin-1', 'replace').decode('latin-1'))
            pdf.set_font("Helvetica", "", 10)
            pdf.multi_cell(0, 5, answer_text.encode('latin-1', 'replace').decode('latin-1'))
            pdf.ln(5)
    return pdf.output(dest='S').encode('latin-1')

def create_roadmap_pdf_bytes(roadmap_df, company_info):
    pdf = PDF(title=f"Strategic Roadmap for {company_info.get('website')}")
    pdf.add_page()
    for _, row in roadmap_df.iterrows():
        pdf.set_font("Helvetica", "B", 14)
        pdf.multi_cell(0, 7, f"Project: {row.get('project_name', 'N/A')}".encode('latin-1', 'replace').decode('latin-1'))
        pdf.ln(2)
        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(40, 7, "Business Value:")
        pdf.set_font("Helvetica", "", 10)
        pdf.cell(0, 7, str(row.get('business_value', 'N/A')))
        pdf.ln()
        pdf.set_font("Helvetica", "B", 10)
        pdf.cell(40, 7, "Level of Effort:")
        pdf.set_font("Helvetica", "", 10)
        pdf.cell(0, 7, str(row.get('level_of_effort', 'N/A')))
        pdf.ln()
        pdf.set_font("Helvetica", "B", 10)
        pdf.multi_cell(0, 5, "Description:")
        pdf.set_font("Helvetica", "", 10)
        pdf.multi_cell(0, 5, str(row.get('description', 'N/A')).encode('latin-1', 'replace').decode('latin-1'))
        pdf.ln(7)
    return pdf.output(dest='S').encode('latin-1')

def generate_powerpoint_bytes(company_info, business_case, competitive_analysis):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Strategic Business Case for {company_info.get('website', 'the Company')}"
    subtitle.text = f"Prepared for: {company_info.get('persona', 'Valued Stakeholder')}\nDate: {datetime.today().strftime('%B %d, %Y')}"
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Business Case & Recommended Strategy"
    tf = body_shape.text_frame
    tf.text = business_case if business_case else "Not generated."
    if competitive_analysis:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = f"Anticipating the {company_info.get('competitor', 'Competitor')} Argument"
        tf = body_shape.text_frame
        tf.text = competitive_analysis
    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream.getvalue()

def format_for_gdocs(export_df, company_info):
    doc_text = [f"## Discovery Q&A for {company_info.get('website')}\n"]
    if 'answer' not in export_df.columns:
        export_df['answer'] = ''
    export_df['answer'] = export_df['answer'].fillna('').astype(str)
    answered_df = export_df[export_df['answer'].str.strip() != ''].copy()
    unanswered_df = export_df[export_df['answer'].str.strip() == ''].copy()
    if not answered_df.empty:
        doc_text.append("\n### Answered Questions\n")
        for category, group in answered_df.groupby('question_category'):
            doc_text.append(f"#### {category}\n")
            for _, row in group.iterrows():
                fav_char = " **(Favorite)**" if row['favorite'] else ""
                doc_text.append(f"**Q: {row['question']}**{fav_char}\n")
                answer = str(row['answer']).replace('\n', '\n> ')
                doc_text.append(f"> {answer}\n")
    if not unanswered_df.empty:
        doc_text.append("\n---\n### Unanswered Questions\n")
        for category, group in unanswered_df.groupby('question_category'):
            doc_text.append(f"#### {category}\n")
            for _, row in group.iterrows():
                doc_text.append(f"- {row['question']}\n")
    return "\n".join(doc_text)

def format_strategy_for_gdocs(initial, refined, competitive):
    doc_text = []
    if initial:
        doc_text.append("## Initial Value Hypothesis\n")
        doc_text.append(initial)
    if refined:
        doc_text.append("\n\n---\n\n## Refined Value Case\n")
        doc_text.append(refined)
    if competitive:
        doc_text.append("\n\n---\n\n## Competitive Battlecard\n")
        doc_text.append(competitive)
    return "\n".join(doc_text)

def create_strategy_csv_bytes(initial, refined, competitive):
    data = []
    if initial:
        data.append({'section': 'Initial Value Hypothesis', 'content': initial})
    if refined:
        data.append({'section': 'Refined Value Case', 'content': refined})
    if competitive:
        data.append({'section': 'Competitive Battlecard', 'content': competitive})
    
    if not data:
        return b''
        
    df = pd.DataFrame(data)
    return df.to_csv(index=False).encode('utf-8')

def format_roadmap_for_gdocs(df):
    if df.empty:
        return "No roadmap has been generated."
    return df.to_markdown(index=False)

def create_briefing_markdown():
    """Compiles all generated content into a single Markdown string."""
    doc_parts = []
    company_name = st.session_state.company_info.get('website', 'Your Target Company')
    doc_parts.append(f"# Sales briefing: {company_name}")

    if st.session_state.company_summary:
        doc_parts.append(f"## Company Overview\n\n{st.session_state.company_summary}")

    initiatives = [cat.split(': ')[1] for cat in st.session_state.questions if cat.startswith("Initiative:")]
    if initiatives:
        doc_parts.append("## Key Initiatives\n\n" + "\n".join([f"- {init}" for init in initiatives]))

    if st.session_state.people_research:
        doc_parts.append("## People Research")
        for person in st.session_state.people_research:
            topics = person.get('topics', {})
            business_topics = ', '.join(topics.get('business', ['N/A']))
            technical_topics = ', '.join(topics.get('technical', ['N/A']))
            personal_topics = ', '.join(topics.get('personal', ['N/A']))
            
            insights_list = "\n".join([f"- {insight}" for insight in person.get('insights', [])])

            person_md = f"""
### {person.get('name', 'N/A')} ({person.get('title', 'N/A')})

**Summary:**
{person.get('summary', 'N/A')}

**Key Insights & Skills:**
{insights_list}

**Suggested Discussion Topics:**
- **Business:** {business_topics}
- **Technical:** {technical_topics}
- **Personal:** {personal_topics}
"""
            doc_parts.append(person_md)

    if st.session_state.initial_value_hypothesis or st.session_state.value_strategy_content:
         doc_parts.append(format_strategy_for_gdocs(st.session_state.initial_value_hypothesis, st.session_state.value_strategy_content, st.session_state.competitive_analysis_content))
    
    if not st.session_state.roadmap_df.empty:
        doc_parts.append("## Strategic Roadmap\n\n" + format_roadmap_for_gdocs(st.session_state.roadmap_df))

    if st.session_state.questions:
        export_data = []
        for category, questions_list in st.session_state.questions.items():
            for question in questions_list:
                export_data.append({"question_category": category, "question": question['text'], "answer": question.get('answer', ''), "favorite": question.get('favorite', False)})
        export_df = pd.DataFrame(export_data)
        doc_parts.append(format_for_gdocs(export_df, st.session_state.company_info))

    if st.session_state.outreach_emails:
        doc_parts.append("## Suggested Outreach Emails")
        for i, (key, email) in enumerate(st.session_state.outreach_emails.items()):
            doc_parts.append(f"### Email Option {i+1}\n\n**Subject:** {email.get('subject')}\n\n**Body:**\n{email.get('body')}")

    return "\n\n---\n\n".join(doc_parts)

def move_question(category, index, direction):
    questions_list = st.session_state.questions[category]
    if direction == 'up' and index > 0:
        questions_list[index], questions_list[index - 1] = questions_list[index - 1], questions_list[index]
    elif direction == 'down' and index < len(questions_list) - 1:
        questions_list[index], questions_list[index + 1] = questions_list[index + 1], questions_list[index]

def toggle_favorite(q_id):
    for category in st.session_state.questions:
        for question in st.session_state.questions[category]:
            if question['id'] == q_id:
                question['favorite'] = not question.get('favorite', False)
                return

def delete_question(q_id):
    for category, q_list in st.session_state.questions.items():
        st.session_state.questions[category] = [q for q in q_list if q['id'] != q_id]

def add_custom_question(category, text_input_key):
    question_text = st.session_state[text_input_key]
    if question_text:
        new_question = {"id": str(uuid.uuid4()), "text": question_text, "answer": "", "favorite": False}
        if category in st.session_state.questions:
            st.session_state.questions[category].append(new_question)
        st.session_state[text_input_key] = ""

def clear_session(preserve_loaded_session_id=False):
    loaded_session_id = st.session_state.get('selected_session_id') if preserve_loaded_session_id else "new"
    current_user = st.session_state.get('current_user')
    selected_model = st.session_state.get('selected_model', 'claude-3-5-sonnet')
    
    keys_to_clear = list(st.session_state.keys())
    for key in keys_to_clear:
        del st.session_state[key]
        
    st.session_state.selected_session_id = loaded_session_id
    st.session_state.selected_model = selected_model
    if current_user:
        st.session_state.current_user = current_user
    
    st.session_state.research_stage = 0
    st.session_state.recommended_initiatives = []
    st.session_state.questions = {}
    st.session_state.company_info = {}
    st.session_state.messages = []
    st.session_state.roadmap_df = pd.DataFrame()
    st.session_state.initial_value_hypothesis = ""
    st.session_state.value_strategy_content = ""
    st.session_state.competitive_analysis_content = ""
    st.session_state.outreach_emails = {}
    st.session_state.notes_content = ""
    st.session_state.company_summary = ""
    st.session_state.people_research = []
    st.session_state.show_briefing = False

# ======================================================================================
# Main App UI
# ======================================================================================
st.sidebar.subheader(f'User: {get_current_user()}')
st.sidebar.subheader("Discovery Session Management")

saved_sessions_df = load_sessions_from_snowflake()
session_options = {"new": "Start a New Session"}
if saved_sessions_df is not None and not saved_sessions_df.empty:
    session_options.update(pd.Series(saved_sessions_df.SESSION_NAME.values, index=saved_sessions_df.SESSION_ID).to_dict())
else:
    st.sidebar.info("ℹ️ No saved sessions found for your user.")

def on_session_change():
    session_id = st.session_state.selected_session_id
    if session_id == "new":
        clear_session()
        st.session_state.show_briefing = False
    else:
        load_selected_session(session_id)
        st.session_state.show_briefing = False

st.sidebar.selectbox(
    "Load a session or start new:",
    options=list(session_options.keys()),
    key='selected_session_id',
    on_change=on_session_change,
    format_func=lambda x: session_options.get(x, "Unknown Session"),
    label_visibility="collapsed"
)

if st.session_state.selected_session_id != "new":
    st.sidebar.button("💾 Update Current Session", on_click=save_session_to_snowflake, use_container_width=True, help="Save changes to the currently loaded session.")
else:
    st.sidebar.button("💾 Save as New Session", on_click=save_session_to_snowflake, use_container_width=True, help="Save the current state as a brand new session.")

if st.session_state.selected_session_id != "new":
    st.sidebar.button("🗑️ Delete Selected Session", use_container_width=True, on_click=delete_session_from_snowflake, args=(st.session_state.selected_session_id,))

st.sidebar.divider()
st.sidebar.subheader("LLM Configuration")
st.sidebar.text("Models impact speed and detail")
model_options = ['claude-3-5-sonnet', 'snowflake-arctic', 'reka-flash', 'reka-core', 'llama3-70b', 'llama3-8b', 'mistral-large', 'gemma-7b']
st.sidebar.selectbox("Select the LLM to use:", options=model_options, key='selected_model', label_visibility="collapsed")
st.sidebar.divider()

if st.sidebar.button("📄 Export Full Briefing", use_container_width=True, help="Compile all session info into a single document."):
    st.session_state.show_briefing = not st.session_state.get('show_briefing', False)
    st.rerun()


if st.session_state.get('token_error_flag'):
    st.warning(
        "🧠 **Context Window Exceeded!** Your last request was too large for the selected model. "
        "Try switching to a model with a larger 'memory' from the sidebar (like `claude-3-5-sonnet` or `reka-core`) and then click the button again.",
        icon="💡"
    )
    st.session_state.token_error_flag = False

if st.session_state.get('show_briefing'):
    st.header("Sales briefing")
    if st.button("⬅️ Back to Main View"):
        st.session_state.show_briefing = False
        st.rerun()
    st.info("Copy the text below and paste it into your document of choice (e.g., Google Docs, Notion).")
    briefing_text = create_briefing_markdown()
    st.text_area("Full briefing", value=briefing_text, height=600)
else:
    tab1, tab2, tab3, tab4, tab5, tab6, tab7 = st.tabs([
        "**Company Research**", "**People Research**", "**Scratchpad**", "**Value & Strategy**", 
        "**Roadmap Builder**", "**Solution Chat**", "**Email Builder**"
    ])

    with tab1:
        col_header1, col_header2 = st.columns([4, 1])
        with col_header2:
            if st.button("Clear & Restart Session", use_container_width=True, help="Clears all data and starts a fresh session."):
                clear_session()
                st.rerun()

        if st.session_state.session_loaded or st.session_state.research_stage == 2:
            with col_header1:
                st.header(f"Company Overview for a {st.session_state.company_info.get('persona', 'Contact')}")
            if st.session_state.company_summary:
                st.info(st.session_state.company_summary)
            st.divider()

            if st.session_state.questions:
                st.header("Step 2: Ask Questions and Capture Responses")
                st.markdown(f"Discovery questions for **{st.session_state.company_info.get('website')}**")

                for cat, q_list in st.session_state.questions.items():
                    for q in q_list:
                        if f"ans_{q['id']}" in st.session_state:
                            q['answer'] = st.session_state[f"ans_{q['id']}"]

                all_categories = list(st.session_state.questions.keys())
                fixed_order = ["Technical Discovery", "Business Discovery"]
                competitive_category = [cat for cat in all_categories if cat.startswith("Competitive Positioning vs.")]
                other_categories = sorted([cat for cat in all_categories if cat not in fixed_order and cat not in competitive_category])
                
                final_category_order = []
                for cat in fixed_order:
                    if cat in all_categories:
                        final_category_order.append(cat)
                final_category_order.extend(competitive_category)
                final_category_order.extend(other_categories)

                for category in final_category_order:
                    questions_list = st.session_state.questions[category]
                    if not questions_list: continue
                    toggle_key = f"toggle_{category.replace(' ', '_')}"
                    is_expanded = st.toggle(f"**{category}** ({len(questions_list)} questions)", key=toggle_key, value=st.session_state.get(toggle_key, False))
                    if is_expanded:
                        with st.container(border=True):
                            for i, question in enumerate(questions_list):
                                q_col1, q_col2, q_col3 = st.columns([0.5, 0.7, 10])
                                with q_col1:
                                    st.button("▲", key=f"up_{question['id']}", on_click=move_question, args=(category, i, 'up'), disabled=(i == 0), help="Move question up")
                                    st.button("▼", key=f"down_{question['id']}", on_click=move_question, args=(category, i, 'down'), disabled=(i == len(questions_list) - 1), help="Move question down")
                                with q_col2:
                                    fav_icon = "⭐" if question.get('favorite', False) else "☆"
                                    st.button(fav_icon, key=f"fav_{question['id']}", on_click=toggle_favorite, args=(question['id'],), help="Toggle Favorite")
                                    st.button("🗑️", key=f"del_{question['id']}", on_click=delete_question, args=(question['id'],), help="Delete Question")
                                with q_col3:
                                    st.markdown(f"**Q: {question['text']}**")
                                    st.text_area("Capture Answer", value=question.get('answer', ''), key=f"ans_{question['id']}", label_visibility="collapsed")
                            st.markdown("---")
                            if st.button(f"Generate 5 More AI Questions", key=f"more_{category.replace(' ', '_')}", use_container_width=True):
                                with st.spinner(f"Generating more questions for {category}..."):
                                    existing_q_texts = [q['text'] for q in questions_list]
                                    info = st.session_state.company_info
                                    new_questions = generate_more_questions_for_category(
                                        info['website'], info['industry'], info['competitor'], info['persona'], 
                                        category, existing_q_texts
                                    )
                                    if new_questions:
                                        st.session_state.questions[category].extend(new_questions)
                                        st.rerun()
                                    else:
                                        st.warning("Could not generate additional questions.")
                            st.markdown("---")
                            st.write("##### Add a Custom Question")
                            custom_q_key = f"custom_q_{category.replace(' ', '_')}"
                            st.text_input("Enter your question:", key=custom_q_key, placeholder="Type your custom question here...", label_visibility="collapsed")
                            st.button("Add Question", key=f"add_q_{category.replace(' ', '_')}", on_click=add_custom_question, args=(category, custom_q_key))
                st.divider()

                st.header("Step 3: Export or Save Responses")
                st.markdown("Once you have captured the answers, you can export them or save them directly to the `SALES_DISCOVERY_ANSWERS` table in Snowflake for broader analytics.")
                export_data = []
                company_name = st.session_state.company_info.get('website', 'N/A').replace('https://','').replace('www.','').split('.')[0].capitalize()
                for category, questions_list in st.session_state.questions.items():
                    for question in questions_list:
                        export_data.append({"id": question['id'], "company_name": company_name, "company_website": st.session_state.company_info.get('website', 'N/A'), "question_category": category, "question": question['text'], "answer": question.get('answer', ''), "favorite": question.get('favorite', False), "saved_at": datetime.utcnow()})
                if export_data:
                    export_df = pd.DataFrame(export_data)
                    col1, col2, col3, col4 = st.columns(4)
                    with col1:
                        st.download_button("📄 Generate PDF", data=create_notes_pdf_bytes(export_df, st.session_state.company_info), file_name=f"discovery_notes_{company_name}.pdf", mime="application/pdf", use_container_width=True)
                    with col2:
                        st.download_button("📥 Download as CSV", data=export_df.to_csv(index=False).encode('utf-8'), file_name=f"discovery_notes_{company_name}.csv", mime='text/csv', use_container_width=True)
                    with col3:
                        if st.button("💾 Save to Answers Table", use_container_width=True, help="Saves all answered questions to the SALES_DISCOVERY_ANSWERS table"):
                            answered_df = export_df[export_df['answer'].str.strip() != ''].copy()
                            if not answered_df.empty:
                                with st.spinner("Saving answered questions..."):
                                    if save_answers_to_snowflake(answered_df):
                                        st.success("✅ Answers saved to Snowflake!")
                            else:
                                st.warning("No answered questions to save.")
                    with col4:
                        if st.button("📋 Copy for Google Docs", use_container_width=True):
                            gdocs_text = format_for_gdocs(export_df, st.session_state.company_info)
                            st.text_area("Formatted Text", value=gdocs_text, height=300)
                            st.info("Use Ctrl+C or Cmd+C to copy the text above.", icon="📋")
        else:
            with col_header1:
                st.header("Step 1: Research Company")
            st.markdown("Enter the company's details to generate a summary and get AI-recommended initiatives.")
            with st.form("research_form"):
                website = st.text_input("**Company Website**", placeholder="e.g., https://www.example.com", value=st.session_state.company_info.get('website', ''))
                industry = st.text_input("**Industry**", placeholder="e.g., Retail, Financial Services", value=st.session_state.company_info.get('industry', ''))
                persona_options = ["C-Level (CEO, CIO, CDO)", "VP of Engineering / Data", "Director of Analytics / BI", "Data Scientist / ML Engineer", "Data/Analytics Engineer"]
                persona = st.selectbox("**Select Title of Person You're Meeting With**", persona_options)
                competitor_options = ["Databricks", "Microsoft Fabric", "AWS", "GCP", "Do Nothing (i.e., keep existing technology)"]
                competitor = st.selectbox("**Select Primary Competitor**", competitor_options)
                research_button = st.form_submit_button("Research Company & Recommend Initiatives", type="primary", disabled=(st.session_state.research_stage > 0))
            if research_button:
                if not website or not industry:
                    st.warning("Please provide both a website and an industry.")
                else:
                    spinner_text = f"Researching company and generating recommendations with {st.session_state.selected_model}..."
                    with st.spinner(spinner_text):
                        st.session_state.company_info = {'website': website, 'industry': industry, 'competitor': competitor, 'persona': persona}
                        st.session_state.company_summary = generate_company_summary(website, industry, persona)
                        st.session_state.recommended_initiatives = recommend_key_initiatives(st.session_state.company_info)
                        st.session_state.research_stage = 1
                        st.rerun()
            if st.session_state.research_stage == 1:
                st.divider()
                if st.session_state.company_summary:
                    st.info(st.session_state.company_summary)
                st.header("Step 2: Confirm Key Initiatives")
                st.markdown("We've recommended the following initiatives based on our research. Please confirm or adjust the selection below, then generate the tailored discovery questions.")
                with st.form("initiatives_form"):
                    standard_initiatives_keys = ["Grow Revenue", "Reduce Risk", "Manage Costs", "Innovate", "Other"]
                    all_options = list(standard_initiatives_keys)
                    for rec in st.session_state.recommended_initiatives:
                        if rec not in all_options:
                            all_options.insert(-1, rec)
                    selected_initiatives = st.multiselect("Key Initiatives", options=all_options, default=st.session_state.recommended_initiatives)
                    other_initiative = st.text_input("**If 'Other', please specify:**")
                    generate_questions_button = st.form_submit_button("Generate Discovery Questions")
                if generate_questions_button:
                    spinner_text = f"Generating tailored discovery questions with {st.session_state.selected_model}..."
                    with st.spinner(spinner_text):
                        final_initiatives = [init for init in selected_initiatives if init != "Other"]
                        if "Other" in selected_initiatives and other_initiative:
                            final_initiatives.append(other_initiative)
                        info = st.session_state.company_info
                        all_questions = generate_discovery_questions(info['website'], info['industry'], info['competitor'], info['persona'])
                        for initiative in final_initiatives:
                            initiative_questions = generate_initiative_questions(info['website'], info['industry'], info['persona'], initiative)
                            if initiative_questions:
                                all_questions[f"Initiative: {initiative}"] = initiative_questions
                            else:
                                st.warning(f"Could not generate questions for the '{initiative}' initiative.")
                        st.session_state.questions = all_questions
                        st.session_state.research_stage = 2
                        save_session_to_snowflake(show_message=False)
                        st.rerun()
    with tab2:
        st.header("People Research")
        st.markdown("Enter an attendee's name and title to generate potential conversation starters.")
        st.info("**Note:** The AI uses its public knowledge base to generate insights based on the person's name, title, and company.", icon="ℹ️")

        if not st.session_state.company_info.get('website'):
            st.warning("Please complete the Company Research on the first tab before researching people.")
        else:
            with st.form("people_form"):
                name = st.text_input("Full Name", placeholder="e.g., Jane Doe")
                title = st.text_input("Title", placeholder="e.g., VP of Engineering")
                research_person_button = st.form_submit_button("Research Attendee")

            if research_person_button:
                if name and title:
                    with st.spinner(f"Researching {name}..."):
                        profile_data = research_person(name, title, st.session_state.company_info)
                        if profile_data:
                            profile_data['name'] = name 
                            profile_data['title'] = title
                            st.session_state.people_research.append(profile_data)
                        else:
                            st.error("Could not generate a profile for this individual.")
                else:
                    st.warning("Please enter both a name and a title.")
            
            st.divider()
            
            if st.session_state.people_research:
                st.subheader("Researched Individuals")
                for i, person in enumerate(st.session_state.people_research):
                    with st.expander(f"**{person.get('name', f'Person {i+1}')}** - {person.get('title')}​"):
                        st.markdown(f"**Summary:**\n{person.get('summary', 'Not available.')}")
                        st.markdown("**Key Insights & Skills:**")
                        for insight in person.get('insights', []):
                            st.markdown(f"- {insight}")
                        
                        st.markdown("**Suggested Discussion Topics:**")
                        topics = person.get('topics', {})
                        st.markdown(f"- **Business:** {', '.join(topics.get('business', ['N/A']))}")
                        st.markdown(f"- **Technical:** {', '.join(topics.get('technical', ['N/A']))}")
                        st.markdown(f"- **Personal:** {', '.join(topics.get('personal', ['N/A']))}")
                        
                        if st.button("🗑️ Delete this profile", key=f"del_person_{i}"):
                            st.session_state.people_research.pop(i)
                            st.rerun()
    with tab3:
        st.header("Freeform Notes & Auto-fill")
        st.markdown("Take your meeting notes here. When you're done, the AI can read your notes and automatically fill out the answers on the 'Company Research' tab.")
        if not st.session_state.get('questions'):
            st.info("Please generate questions on the **'Company Research'** tab first.")
        else:
            notes = st.text_area("Meeting Notes", key='notes_content', placeholder="Start typing your meeting notes here...", height=400)
            if st.button("📝 Auto-fill Answers from Notes", type="primary"):
                if notes:
                    spinner_text = f"Analyzing notes with {st.session_state.selected_model}..."
                    with st.spinner(spinner_text):
                        populated_questions = autofill_answers_from_notes(notes, st.session_state.questions)
                        if populated_questions:
                            st.session_state.questions = populated_questions
                            st.success("Auto-fill complete! Check the 'Company Research' tab to see the results.")
                        else:
                            st.error("The AI could not process the notes to fill out the questions. Please try again.")
                else:
                    st.warning("Please enter some notes before trying to auto-fill.")
    
    with tab4:
        st.header("Value & Strategy")
        st.markdown("Generate an initial value hypothesis based on public info, then refine it with captured discovery notes.")
        if not st.session_state.company_info.get('website'):
            st.info("Please complete the research on the **'Company Research'** tab first.")
        else:
            with st.container(border=True):
                st.subheader("Initial Value Hypothesis (Pre-Discovery)")
                if st.button("💡 Generate Initial Hypothesis", help="Generates a hypothesis based on company website and industry.", disabled=bool(st.session_state.initial_value_hypothesis)):
                    spinner_text = f"Generating initial hypothesis with {st.session_state.selected_model}..."
                    with st.spinner(spinner_text):
                        hypothesis = generate_initial_value_hypothesis(st.session_state.company_info)
                        st.session_state.initial_value_hypothesis = hypothesis
                        st.rerun()
                if st.session_state.initial_value_hypothesis:
                    st.markdown(st.session_state.initial_value_hypothesis)
            st.divider()
            with st.container(border=True):
                st.subheader("Refined Value Case (Post-Discovery)")
                notes_list = []
                if st.session_state.get('questions'):
                    for category, questions_list in st.session_state.questions.items():
                        for q in questions_list:
                            answer = str(q.get('answer', '')).strip()
                            if answer: notes_list.append(f"Category: {category}\nQ: {q['text']}\nA: {answer}\n")
                discovery_notes_str = "\n".join(notes_list)
                if not discovery_notes_str:
                    st.warning("Please answer at least one question on the first tab to generate a refined value case.")
                else:
                    if st.button("📈 Generate Refined Business Case & Strategy", type="primary"):
                        spinner_text = f"Building refined case with {st.session_state.selected_model}..."
                        with st.spinner(spinner_text):
                            business_case = generate_business_case(st.session_state.company_info, discovery_notes_str)
                            st.session_state.value_strategy_content = business_case
                            st.session_state.competitive_analysis_content = ""
                    if st.session_state.value_strategy_content:
                        st.markdown(st.session_state.value_strategy_content)
                        st.divider()
                        st.subheader("Competitive Battlecard")
                        button_text = f"🛡️ Generate 'Steel Man' Argument for {st.session_state.company_info.get('competitor')}"
                        if st.button(button_text):
                            spinner_text = f"Preparing counter-arguments with {st.session_state.selected_model}..."
                            with st.spinner(spinner_text):
                                competitive_arg = generate_competitive_argument(st.session_state.company_info, discovery_notes_str)
                                st.session_state.competitive_analysis_content = competitive_arg
                        if st.session_state.competitive_analysis_content:
                            st.markdown(st.session_state.competitive_analysis_content)
            
            st.divider()
            if st.session_state.initial_value_hypothesis or st.session_state.value_strategy_content:
                st.subheader("Export Value & Strategy")
                strat_col1, strat_col2, strat_col3 = st.columns(3)
                with strat_col1:
                    ppt_bytes = generate_powerpoint_bytes(st.session_state.company_info, st.session_state.value_strategy_content, st.session_state.competitive_analysis_content)
                    st.download_button(label="📊 Download as PowerPoint", data=ppt_bytes, file_name=f"Strategy_{st.session_state.company_info.get('website', 'export').replace('https://','').replace('www.','').split('.')[0]}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
                with strat_col2:
                    csv_data = create_strategy_csv_bytes(st.session_state.initial_value_hypothesis, st.session_state.value_strategy_content, st.session_state.competitive_analysis_content)
                    st.download_button(label="📥 Download as CSV", data=csv_data, file_name=f"strategy_{st.session_state.company_info.get('website', 'export').replace('https://','').replace('www.','').split('.')[0]}.csv", mime="text/csv", use_container_width=True)
                with strat_col3:
                    if st.button("📋 Copy for Google Docs", key="gdocs_strategy", use_container_width=True):
                        gdocs_text = format_strategy_for_gdocs(st.session_state.initial_value_hypothesis, st.session_state.value_strategy_content, st.session_state.competitive_analysis_content)
                        st.text_area("Formatted Text for Google Docs", value=gdocs_text, height=300)
                        st.info("Use Ctrl+C or Cmd+C to copy the text above.", icon="📋")

    with tab5:
        st.header("Roadmap Builder")
        if not st.session_state.get('questions'):
            st.info("Please complete the research on the **'Company Research'** tab first.")
        else:
            notes_list = []
            for category, questions_list in st.session_state.questions.items():
                for q in questions_list:
                    answer = str(q.get('answer', '')).strip()
                    if answer: notes_list.append(f"Category: {category}\nQ: {q['text']}\nA: {answer}\n")
            discovery_notes_str = "\n".join(notes_list)
            if not discovery_notes_str:
                st.warning("Please answer at least one question to generate a roadmap.")
            else:
                st.markdown("Based on the captured discovery answers, generate a potential implementation plan.")
                priority = st.radio("**Prioritize roadmap by:**", ("Quick Wins (Lowest Effort First)", "Highest Business Value First"), horizontal=True)
                if st.button("🚀 Generate Strategic Roadmap", type="primary"):
                    spinner_text = f"Building the roadmap with {st.session_state.selected_model}..."
                    with st.spinner(spinner_text):
                        st.session_state.roadmap_df = generate_roadmap(st.session_state.company_info, discovery_notes_str, priority)
        if not st.session_state.roadmap_df.empty:
            st.divider()
            st.subheader("Generated Roadmap")
            st.dataframe(st.session_state.roadmap_df, use_container_width=True)
            st.subheader("Export Roadmap")
            roadmap_col1, roadmap_col2, roadmap_col3 = st.columns(3)
            company_name = st.session_state.company_info.get('website', 'N/A').replace('https://','').replace('www.','').split('.')[0].capitalize()
            with roadmap_col1:
                pdf_bytes = create_roadmap_pdf_bytes(st.session_state.roadmap_df, st.session_state.company_info)
                st.download_button(label="📄 Download Roadmap as PDF", data=pdf_bytes, file_name=f"roadmap_{company_name}.pdf", mime="application/pdf", use_container_width=True)
            with roadmap_col2:
                csv_data = st.session_state.roadmap_df.to_csv(index=False).encode('utf-8')
                st.download_button(label="📥 Download Roadmap as CSV", data=csv_data, file_name=f"roadmap_{company_name}.csv", mime="text/csv", use_container_width=True)
            with roadmap_col3:
                 if st.button("📋 Copy for Google Docs", key="gdocs_roadmap", use_container_width=True):
                    gdocs_text = format_roadmap_for_gdocs(st.session_state.roadmap_df)
                    st.text_area("Formatted Text for Google Docs", value=gdocs_text, height=250)
                    st.info("Use Ctrl+C or Cmd+C to copy the text above.", icon="📋")

    with tab6:
        st.header("Snowflake Solution Chat")
        if not st.session_state.company_info.get('website'):
            st.info("Please research an account on the first tab to activate the chat.")
        else:
            st.markdown(f"Ask questions about how Snowflake can help **{st.session_state.company_info.get('website')}**.")
            for message in st.session_state.messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
            if prompt := st.chat_input("How can Snowflake help with their data challenges?"):
                st.session_state.messages.append({"role": "user", "content": prompt})
                with st.chat_message("user"):
                    st.markdown(prompt)
                with st.chat_message("assistant"):
                    message_placeholder = st.empty()
                    spinner_text = f"Thinking with {st.session_state.selected_model}..."
                    with st.spinner(spinner_text):
                        full_response = get_chatbot_response(st.session_state.company_info, st.session_state.messages)
                        message_placeholder.markdown(full_response)
                    st.session_state.messages.append({"role": "assistant", "content": full_response})

    with tab7:
        st.header("Outreach Email Generator")
        if 'roadmap_df' not in st.session_state or st.session_state.roadmap_df.empty:
            st.info("Please generate a roadmap on the **'Roadmap Builder'** tab first to enable email generation.")
        else:
            st.markdown("Draft personalized follow-up emails based on the discovery notes and the strategic roadmap.")
            def handle_regenerate(email_key_to_replace):
                notes_list = []
                for category, questions_list in st.session_state.questions.items():
                    for q in questions_list:
                        answer = str(q.get('answer', '')).strip()
                        if answer: notes_list.append(f"Category: {category}\nQ: {q['text']}\nA: {answer}\n")
                discovery_notes_str = "\n".join(notes_list)
                spinner_text = f"Re-generating email with {st.session_state.selected_model}..."
                with st.spinner(spinner_text):
                    new_email_data = regenerate_single_email(
                        st.session_state.company_info, discovery_notes_str, st.session_state.roadmap_df, 
                        st.session_state.outreach_emails, email_key_to_replace
                    )
                    if new_email_data and "subject" in new_email_data and "body" in new_email_data:
                        st.session_state.outreach_emails[email_key_to_replace] = new_email_data
                    else:
                        st.warning("Could not re-generate the email. Please try again.")
            if st.button("✉️ Draft Follow-up Emails", type="primary"):
                notes_list = []
                for category, questions_list in st.session_state.questions.items():
                    for q in questions_list:
                        answer = str(q.get('answer', '')).strip()
                        if answer: notes_list.append(f"Category: {category}\nQ: {q['text']}\nA: {answer}\n")
                discovery_notes_str = "\n".join(notes_list)
                spinner_text = f"Drafting emails with {st.session_state.selected_model}..."
                with st.spinner(spinner_text):
                    st.session_state.outreach_emails = generate_outreach_emails(st.session_state.company_info, discovery_notes_str, st.session_state.roadmap_df)
            if st.session_state.outreach_emails:
                st.divider()
                for i, (email_key, email_data) in enumerate(st.session_state.outreach_emails.items()):
                    subject = email_data.get('subject', "No Subject")
                    expander_title = f"**Email Version {i+1}: {subject}**"
                    with st.expander(expander_title, expanded=(i==0)):
                        b_col1, b_col2 = st.columns([4, 1])
                        with b_col1:
                            st.markdown(f"##### Subject: {subject}")
                        with b_col2:
                            st.button("🔄 Re-generate", key=f"regen_{email_key}", on_click=handle_regenerate, args=(email_key,))
                        st.markdown("---")
                        body = email_data.get('body', "Could not generate email body.").replace('\n', '<br>')
                        st.markdown(body, unsafe_allow_html=True)
                        st.markdown("---")
                        full_email_text = f"Subject: {subject}\n\n{email_data.get('body', '')}"
                        st.text_area("**Copyable Email Content**", full_email_text, height=200, key=f"copy_text_{email_key}")
